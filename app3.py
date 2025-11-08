import os
import json
from flask import Flask, request, jsonify, render_template_string

# --- NEW/UPDATED IMPORTS ---
try:
    from pytube import YouTube # For downloading YouTube audio
    PYTUBE_AVAILABLE = True
except ImportError:
    PYTUBE_AVAILABLE = False
# ---------------------------

from googletrans import Translator, LANGUAGES
import fitz # PyMuPDF for PDF
from docx import Document # python-docx
from pptx import Presentation # python-pptx
import speech_recognition as sr
import tempfile
from datetime import datetime
import subprocess 
from difflib import get_close_matches # For filename suggestion
import urllib.parse # For URL encoding
import re # For parsing YouTube search results

# Attempt to import requests for YouTube search/web actions
try:
    import requests
    REQUESTS_AVAILABLE = True
except ImportError:
    REQUESTS_AVAILABLE = False

# Attempt to import moviepy for video processing; essential for transcription
try:
    from moviepy import VideoFileClip
    MOVIEPY_AVAILABLE = True
except ImportError:
    MOVIEPY_AVAILABLE = False


# --- CONFIGURATION & GLOBAL STATE (Server-Side) ---

# Set the file path to the user's Downloads folder
FILE_STORAGE_PATH = os.path.join(os.path.expanduser("~"), "Downloads") 

SUPPORTED_EXTENSIONS = ('.txt', '.pdf', '.docx', '.pptx', '.mp4', '.mov', '.avi', '.mkv', '.wmv')

os.makedirs(FILE_STORAGE_PATH, exist_ok=True) # Ensure the folder exists

app = Flask(__name__)
translator = Translator()
recognizer = sr.Recognizer()

# Initial State Defaults
INITIAL_FILE_CONTENT = ""
INITIAL_READ_POSITION = 0
INITIAL_TRANSLATION_LANGUAGE = 'en' # Target language code

# Global State Variables
current_file_content = INITIAL_FILE_CONTENT # Stores the original content for file reading/transcription
read_position = INITIAL_READ_POSITION
translation_language = INITIAL_TRANSLATION_LANGUAGE # Target language code


# --- HELPER FUNCTION: RESET STATE ---

def handle_reset_command():
    """Resets the global state variables to their initial values."""
    global current_file_content, read_position, translation_language
    
    current_file_content = INITIAL_FILE_CONTENT
    read_position = INITIAL_READ_POSITION
    
    return {
        "status": "success", 
        "action": "stop_read", # Instruct client to stop any ongoing reading
        "message": "Assistant memory cleared. Current file closed, and reading stopped. Ready for the next command.",
        "has_content": False
    }

# --- HELPER FUNCTIONS ---

def get_language_code(language_name):
    """Converts a language name (e.g., 'French') to a code (e.g., 'fr')."""
    for code, name in LANGUAGES.items():
        if name.lower() == language_name.lower():
            return code
    return None

def extract_language_codes(command):
    """
    Parses the command for source and destination languages for transcription/translation.
    """
    global translation_language

    command = command.lower()
    target_code = None
    source_code = None

    # Pre-process command to remove the filename to clean up language parsing
    clean_command = command
    for ext in SUPPORTED_EXTENSIONS:
        if ext in command:
            clean_command = command.split(ext, 1)[0].strip()
            break
            
    # 1. Extract Target Language (translate to [language])
    if "translate to" in clean_command:
        parts = clean_command.split("translate to", 1)[1].strip().split()
        target_name_parts = []
        for part in parts:
            if part == 'from':
                break
            target_name_parts.append(part)
            
        target_name = " ".join(target_name_parts)
        target_code = get_language_code(target_name)

    # 2. Extract Source Language (from [language])
    if "from" in clean_command:
        parts = clean_command.split("from", 1)[1].strip().split()
        if parts:
            for i in range(min(3, len(parts))): # Check 1, 2, or 3 words
                lang_candidate = " ".join(parts[:i+1])
                code = get_language_code(lang_candidate)
                if code:
                    source_code = code
                    break

    # 3. Apply Defaults (Prioritizing Global State for Target)
    if target_code is None and ("transcribe" in command or "translate" in command):
        target_code = translation_language  
    
    # 4. Handle STT regional codes (Hindi/Hinglish often needs the regional code)
    if source_code == 'hi' or ("transcribe" in command and source_code is None):
        if 'hindi' in command or 'hinglish' in command:
             source_code = 'hi-IN' 

    # If the user is transcribing and didn't specify source, default to 'en'
    if source_code is None and "transcribe" in command:
        source_code = 'en' # Fallback for transcription
        
    return target_code, source_code

def extract_filename_from_command(command):
    """
    Parses the command to find a file name, handling common STT errors.
    """
    command = command.lower()
    
    # Replace spoken "dot" forms with a literal dot 
    command = command.replace(" dot ", ".")
    command = command.replace(" point ", ".")
    command = command.replace(" mp 4", " mp4")

    # Expanded keywords and extensions to improve parsing
    file_keywords = ["open", "read", "play", "transcribe", "file", "document", "video", "docx", "pdf", "pptx", "txt", "mp4", "mov", "avi", "mkv", "wmv"]
    
    for keyword in file_keywords:
        if keyword in command:
            parts = command.split(keyword, 1)
            if len(parts) > 1:
                filename_candidate = parts[1].strip()
                # Clean up everything after the file name (like "and translate...")
                for separator in [' and ', ' to ', ' in ', ' from ']:
                    if separator in filename_candidate:
                        filename_candidate = filename_candidate.split(separator)[0]
                
                # Check if it looks like a file name
                if "." in filename_candidate and len(filename_candidate) > 3:
                    final_filename = filename_candidate.strip().rstrip('.')
                    return final_filename
    return None

def extract_file_content(full_path):
    """Extracts content based on file extension."""
    content = ""
    file_name = os.path.basename(full_path)
    file_extension = os.path.splitext(file_name)[1].lower()

    if file_extension == ".txt":
        with open(full_path, "r", encoding="utf-8") as file:
            content = file.read()
            
    elif file_extension == ".docx":
        doc = Document(full_path)
        content = "\n".join([para.text for para in doc.paragraphs])
        
    elif file_extension == ".pdf":
        doc = fitz.open(full_path)
        content = "\n".join(page.get_text() for page in doc)
        doc.close()
        
    elif file_extension == ".pptx":
        prs = Presentation(full_path)
        content = "\n".join(
            [
                "\n".join(
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ) for slide in prs.slides
            ]
        )
    return content

# NEW: Filename Suggestion Helper
def suggest_filename(wrong_name, file_list):
    """Suggests close matches from a list of files using difflib."""
    suggestions = get_close_matches(wrong_name, file_list, n=3, cutoff=0.6)
    return suggestions

def transcribe_video_audio(file_name, requested_target_lang='en', requested_source_lang='en'):
    """
    Extracts audio from a local video, transcribes it, and translates the result.
    
    *** MODIFIED: Uses requested_source_lang for STT recognition. ***
    """
    global current_file_content, read_position, translation_language

    if not MOVIEPY_AVAILABLE:
        return {"status": "error", "message": "Video transcription failed: **moviepy** library is not installed on the server."}
    
    full_path = os.path.join(FILE_STORAGE_PATH, file_name)

    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmpfile:
        audio_file_path = tmpfile.name

    try:
        # 1. EXTRACT AUDIO
        clip = VideoFileClip(full_path)
        clip.audio.write_audiofile(audio_file_path, logger=None)
        clip.close()
        
        # --- Use the requested source language for STT ---
        stt_lang_code = requested_source_lang
        
        # 2. TRANSCRIBE AUDIO
        with sr.AudioFile(audio_file_path) as source:
            audio = recognizer.record(source)  
            transcribed_text = recognizer.recognize_google(audio, language=stt_lang_code)  # <--- LANGUAGE CODE APPLIED
            
            if not transcribed_text:
                source_name = LANGUAGES.get(stt_lang_code, stt_lang_code)
                raise sr.UnknownValueError(f"No speech detected after trying {source_name} transcription.")

            # --- Dynamic Translation ---
            detection_result = translator.detect(transcribed_text)
            actual_source_lang_code = detection_result.lang
            actual_source_lang_name = LANGUAGES.get(actual_source_lang_code, 'Unknown Language')
            
            target_lang_name = LANGUAGES.get(requested_target_lang, requested_target_lang)
            
            final_content = transcribed_text
            translation_message = ""
            
            if actual_source_lang_code != requested_target_lang:
                translation_result = translator.translate(transcribed_text, dest=requested_target_lang)
                final_content = translation_result.text
                translation_message = f" (Translated from **{actual_source_lang_name}** to **{target_lang_name}**.)"
            else:
                translation_message = f" (Language detected as **{actual_source_lang_name}**. No translation needed.)"

    except sr.UnknownValueError as e:
        return {"status": "error", "message": f"Transcription failed: Could not understand the audio in the video file. Error: {e}"}
    except sr.RequestError as e:
        return {"status": "error", "message": f"Transcription failed: Could not reach Google Speech Recognition service ({e})."}
    except Exception as e:
        print(f"ERROR: Video transcription failed for {file_name}: {e}")
        return {"status": "error", "message": f"Error during video processing or transcription: {e}"}
        
    finally:
        # CLEANUP
        if os.path.exists(audio_file_path):
            os.remove(audio_file_path)

    if not final_content:
        return {"status": "error", "message": f"Successfully processed the file, but no readable script was generated."}

    # Set global state for reading
    current_file_content = final_content 
    read_position = 0 
    translation_language = requested_target_lang # Set reading language to the requested target
    
    return {
        "status": "success", 
        "action": "start_read",
        "message": f"Transcription complete. Starting to read the script for '{file_name}' in **{target_lang_name}**.{translation_message}"
    }

# NEW: REAL YOUTUBE TRANSCRIPTION WITH PYTUBE ---

def transcribe_youtube_audio_real(query, requested_target_lang, requested_source_lang='en'):
    """
    Downloads YouTube audio using pytube, transcribes it, and translates the result.
    
    *** MODIFIED: Now accepts and uses requested_source_lang for STT recognition. ***
    """
    global current_file_content, read_position, translation_language

    if not PYTUBE_AVAILABLE:
        return {"status": "error", "message": "YouTube transcription failed: The **pytube** library is not imported. Please run `pip install pytube`."}
    
    if not REQUESTS_AVAILABLE:
        return {"status": "error", "message": "YouTube transcription requires the **requests** library for searching. Please install it with `pip install requests`."}
        
    # 1. Get Video ID (using existing search function)
    video_id = search_youtube_simple(query)
    if not video_id:
           return {"status": "error", "message": f"Could not find a YouTube video for the query: '{query}'."}
    
    video_url = f"https://www.youtube.com/watch?v={video_id}"
    temp_dir = tempfile.gettempdir()
    downloaded_file = None
    
    try:
        # 2. Download Audio using pytube
        yt = YouTube(video_url)
        audio_stream = yt.streams.filter(only_audio=True).first()
        
        if not audio_stream:
            return {"status": "error", "message": "Could not find an audio-only stream for the selected YouTube video."}

        # Temporary filename with original extension (.mp4 or .webm)
        audio_filename = f"yt_audio_{video_id}{audio_stream.subtype}"
        downloaded_file = audio_stream.download(output_path=temp_dir, filename=audio_filename)

        if not downloaded_file or not os.path.exists(downloaded_file):
            raise Exception("Pytube download failed to produce a file.")
            
        print(f"Downloaded audio to: {downloaded_file}")

        # 3. Transcribe Audio (using SpeechRecognition)
        # --- Use the requested source language for STT ---
        stt_lang_code = requested_source_lang
        
        with sr.AudioFile(downloaded_file) as source:
            recognizer.adjust_for_ambient_noise(source)
            audio = recognizer.record(source)  
            transcribed_text = recognizer.recognize_google(audio, language=stt_lang_code)  # <--- LANGUAGE CODE APPLIED
            
            if not transcribed_text:
                raise sr.UnknownValueError("No speech detected in the downloaded YouTube audio.")
                
            # 4. Handle Translation
            detection_result = translator.detect(transcribed_text)
            actual_source_lang_code = detection_result.lang
            actual_source_lang_name = LANGUAGES.get(actual_source_lang_code, 'Unknown Language')
            target_lang_name = LANGUAGES.get(requested_target_lang, requested_target_lang)
            
            final_content = transcribed_text
            translation_message = ""
            
            if actual_source_lang_code != requested_target_lang:
                translation_result = translator.translate(transcribed_text, dest=requested_target_lang)
                final_content = translation_result.text
                translation_message = f" (Translated from **{actual_source_lang_name}** to **{target_lang_name}**.)"
            else:
                translation_message = f" (Language detected as **{actual_source_lang_name}**. No translation needed.)"

    except sr.UnknownValueError as e:
        return {"status": "error", "message": f"Transcription failed: Could not understand the audio. ({e})"}
    except Exception as e:
        return {"status": "error", "message": f"YouTube download/transcription error: {e}"}
        
    finally:
        # CLEANUP TEMPORARY FILE
        if downloaded_file and os.path.exists(downloaded_file):
            os.remove(downloaded_file)

    # 5. Update Global State and Return
    current_file_content = final_content 
    read_position = 0 
    translation_language = requested_target_lang
    
    return {
        "status": "success", 
        "action": "start_read",
        "message": f"Transcription for YouTube video '{yt.title}' complete. Starting to read the script in **{target_lang_name}**." + translation_message
    }


# --- YOUTUBE SEARCH FUNCTIONS (Kept clean) ---

def search_youtube_simple(query):
    """
    Simple YouTube search using requests - returns the first video ID or None.
    """
    if not REQUESTS_AVAILABLE:
        return None
    
    try:
        encoded_query = urllib.parse.quote(query)
        search_url = f"https://www.youtube.com/results?search_query={encoded_query}"
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(search_url, headers=headers, timeout=10)
        
        if response.status_code == 200:
            pattern = r'"videoId":"([a-zA-Z0-9_-]{11})"|/watch\?v=([a-zA-Z0-9_-]{11})'
            matches = re.findall(pattern, response.text)
            
            if matches:
                for match_tuple in matches:
                    video_id = next((id for id in match_tuple if id), None) 
                    if video_id:
                        return video_id 
            return None
            
    except Exception as e:
        print(f"YouTube search error: {e}")
        pass
    
    return None

def extract_youtube_query(command):
    """Extracts the search query from YouTube commands."""
    command = command.lower().strip()
    
    triggers = ["play ", "search for ", "search ", "find ", "show ", "watch ", "transcribe ", "translate "]
    
    for trigger in triggers:
        if command.startswith(trigger):
            command = command[len(trigger):].strip()
            break
    
    if command.endswith("on youtube"):
        command = command[:-10].strip()
    elif command.endswith("in youtube"):
        command = command[:-10].strip()
    elif command.endswith("youtube"):
        command = command[:-7].strip()
        
    for lang_code, lang_name in LANGUAGES.items():
        if f"to {lang_name.lower()}" in command:
            command = command.replace(f"to {lang_name.lower()}", "").strip()
        if f"from {lang_name.lower()}" in command:
            command = command.replace(f"from {lang_name.lower()}", "").strip()

    command = command.replace(" and ", " ").replace(" with ", " ").strip()
    
    return " ".join(command.split()) if command else None

def handle_youtube_command(command):
    """
    Handles YouTube play commands using simple search.
    """
    query = extract_youtube_query(command)
    
    if not query:
        return None
    
    if not REQUESTS_AVAILABLE:
        return {"status": "error", "message": "YouTube search and playback feature requires the **requests** library. Please install it with: `pip install requests`"}
    
    # Try to get video ID
    video_id = search_youtube_simple(query)
    
    if video_id:
        return {
            "status": "action",
            "action": "play_youtube_embedded", 
            "video_id": video_id,
            "query": query,
            "video_title": query, 
            "message": f"Playing '{query}' on YouTube."
        }
    
    # Fallback: open search page
    encoded_query = urllib.parse.quote(query)
    youtube_search_url = f"https://www.youtube.com/results?search_query={encoded_query}"
    
    return {
        "status": "action",
        "action": "play_youtube",
        "url": youtube_search_url,
        "query": query,
        "encoded_query": encoded_query,
        "message": f"Opening YouTube search for '{query}'."
    }

# --- CORE LOGIC HANDLERS (Continued) ---

def handle_open_file(file_name_raw, command_text):
    """Opens a file, extracts content, and resets the reading state."""
    global current_file_content, read_position, translation_language

    file_name = file_name_raw.strip()
    full_path = os.path.join(FILE_STORAGE_PATH, file_name)

    if not os.path.exists(full_path):
        # *** Filename Suggestion Logic ***
        all_files = [
            f for f in os.listdir(FILE_STORAGE_PATH)
            if os.path.isfile(os.path.join(FILE_STORAGE_PATH, f)) and f.lower().endswith(SUPPORTED_EXTENSIONS)
        ]
        
        suggestions = suggest_filename(file_name, all_files)
        
        error_message = f"The file '{file_name}' does not exist in the server's Downloads folder."
        if suggestions:
            suggestion_list = ", ".join([f"'{s}'" for s in suggestions])
            error_message += f" Did you mean: **{suggestion_list}**?"
            
        return {"status": "error", "message": error_message}
        
    file_extension = os.path.splitext(file_name)[1].lower()
    
    # --- File Groups for Correct Handling ---
    TEXT_EXTRACTION_FILES = ['.txt', '.pdf', '.docx', '.pptx'] 
    MEDIA_FILES = ['.mp4', '.mov', '.avi', '.mkv', '.wmv'] 

    # 1. Handle Transcription (If explicitly asked for media)
    if file_extension in MEDIA_FILES and "transcribe" in command_text:
        # --- MODIFIED: Extract BOTH source and target languages ---
        target_lang, source_lang = extract_language_codes(command_text)
        
        target_lang = target_lang if target_lang else translation_language  
        source_lang = source_lang if source_lang else 'en'
        
        return transcribe_video_audio(file_name, target_lang, source_lang)
    
    # 2. Handle System Open (For media files only, if transcription not requested)
    if file_extension in MEDIA_FILES:
        try:
            if os.name == 'nt': # Windows
                os.startfile(full_path) 
            elif os.name == 'posix': # Linux/macOS
                subprocess.call(['xdg-open', full_path]) # Linux/Unix
            else:
                return {"status": "info", "message": f"Cannot automatically open media file '{file_name}' on this OS."}
            
            return {
                "status": "info", 
                "message": f"Opening media file '{file_name}' using your default system application."
            }
        except Exception as e:
            return {"status": "error", "message": f"Could not open media file locally. Error: {e}"}

    # 3. Handle Standard Reading (For all document types)
    if file_extension in TEXT_EXTRACTION_FILES:
        try:
            content = extract_file_content(full_path)
            
            if content:
                current_file_content = content
                read_position = 0
                
                return {
                    "status": "success", 
                    "action": "start_read", 
                    "message": f"File '{file_name}' loaded successfully ({len(content)} characters). Starting to read in {LANGUAGES.get(translation_language, 'English')}."
                }
            
            return {"status": "error", "message": f"The file {file_name} was opened, but no readable content was extracted."}

        except Exception as e:
            return {"status": "error", "message": f"Error opening/reading file: {e}"}

    # 4. Fallback
    return {"status": "error", "message": f"File type {file_extension} is not supported for reading or local opening."}


def get_next_chunk(chunk_size=500):
    """Fetches the next chunk of content, translates it if necessary."""
    global read_position, current_file_content, translation_language

    if not current_file_content:
        return {"status": "done", "chunk": "", "message": "No file loaded."}

    if read_position >= len(current_file_content):
        read_position = len(current_file_content) 
        return {"status": "done", "chunk": "", "message": "End of file reached."}

    chunk_end = min(read_position + chunk_size, len(current_file_content))
    original_chunk = current_file_content[read_position:chunk_end]

    translated_chunk = original_chunk
    if translation_language != 'en':
        try:
            translated_chunk = translator.translate(original_chunk, dest=translation_language).text
        except Exception as e:
            pass 

    read_position = chunk_end 
    
    return {
        "status": "reading",
        "chunk": translated_chunk,
        "message": f"Reading chunk in {LANGUAGES.get(translation_language, translation_language)}.",
        "progress": int((read_position / len(current_file_content)) * 100)
    }

def handle_translation_command(command):
    """Changes the target translation language."""
    global translation_language
    
    parts = command.split('to')
    if len(parts) > 1:
        lang_name = parts[-1].strip()
        lang_code = get_language_code(lang_name)
        
        if lang_code:
            translation_language = lang_code
            return {"status": "success", "message": f"Translation language set to {LANGUAGES[lang_code]}."}
        
        return {"status": "error", "message": f"Sorry, I don't recognize the language '{lang_name}'."}
    
    return {"status": "info", "message": f"The current translation language is {LANGUAGES.get(translation_language, translation_language)}."}

def handle_control_command(command):
    """Handles commands like pause, resume, restart, stop, and reading speed."""
    global read_position, current_file_content
    
    # Check for reading rate command
    if command.startswith("set reading speed to"):
        try:
            rate_str = command.split("to")[-1].strip()
            rate = float(rate_str)
            # Send the rate back to the client via a new action type
            return {"status": "action", "action": "set_rate", "rate": rate, "message": f"Reading speed set to {rate}x."}
        except ValueError:
            return {"status": "error", "message": "Invalid reading speed provided. Please use a number (e.g., set reading speed to 1.5)."}

    if command == "restart":
        if not current_file_content:
            return {"status": "info", "message": "I don't have a file loaded to restart reading."}
        read_position = 0
        return {"status": "action", "action": "start_read", "message": f"Restarted reading session from the beginning in {LANGUAGES.get(translation_language, translation_language)}."}
        
    elif command == "resume":
        if not current_file_content:
            return {"status": "info", "message": "No file content loaded to resume."}
        if read_position >= len(current_file_content):
            return {"status": "info", "message": "The file has already finished reading. Use 'restart' to begin again."}
        return {"status": "action", "action": "start_read", "message": f"Resuming reading from the previous position in {LANGUAGES.get(translation_language, translation_language)}."}

    elif command in ("stop", "pause"):
        return {"status": "info", "message": "Reading paused. Say 'resume' to continue."}
    
    return {"status": "info", "message": "Command recognized, but no immediate server action needed."}

def handle_web_action_command(question):
    """Handles commands that result in opening a web page on the client."""
    
    if "order" in question:
        item = question.replace("order", "").replace("online", "").strip()
        
        if "pizza" in item:
            url = "https://www.dominos.com/"
            message = "Redirecting you to Domino's to order pizza."
        elif "book" in item:
            url = "https://www.amazon.com/"
            message = "Redirecting you to Amazon to order a book."
        elif item:
            url = f"https://www.google.com/search?q=best+quality+{item}+to+order"
            message = f"I've searched for the best quality {item}."
        else:
            return {"status": "info", "message": "Please specify what you want to order."}
        
        return {"status": "action", "action": "open_url", "url": url, "message": message}

    elif "search" in question:
        query = question.replace("search", "").replace("web", "").strip()
        if query:
            search_url = f"https://www.google.com/search?q={query}"
            return {"status": "action", "action": "open_url", "url": search_url, "message": f"I've searched the web for: {query}."}
        else:
            return {"status": "info", "message": "Please specify what you want to search for."}
            
    return None

# --- NEW: Endpoint to list files in the storage directory ---
@app.route('/list_files')
def list_files():
    """Returns a list of readable and video files in the storage path."""
    
    try:
        files = [
            f for f in os.listdir(FILE_STORAGE_PATH)
            if os.path.isfile(os.path.join(FILE_STORAGE_PATH, f)) and f.lower().endswith(SUPPORTED_EXTENSIONS)
        ]
        return jsonify({"status": "success", "files": files})
    except Exception as e:
        return jsonify({"status": "error", "message": f"Failed to list files: {e}"}), 500


# --- FLASK ROUTES ---

@app.route('/')
def index():
    """Serves the main HTML/JS front-end page."""
    global translation_language
    
    # Generate the list of supported languages for the dropdown
    language_options = "".join([
        f'<option value="{code}" {"selected" if code == translation_language else ""}>{name}</option>'
        for code, name in LANGUAGES.items()
    ])
    
    # Check if requests is available to warn user if YouTube autoplay is disabled
    youtube_status_message = ""
    if not REQUESTS_AVAILABLE:
        youtube_status_message = """
        <div class="bg-yellow-100 border-l-4 border-yellow-500 text-yellow-700 p-3 mb-4 rounded-lg shadow-md" role="alert">
            <p class="font-bold">YouTube Feature Warning</p>
            <p class="text-sm">The automated **YouTube search and play/transcribe** features require the <code>requests</code> library. Please install it with: <code>pip install requests</code></p>
        </div>
        """
    
    transcribe_deps_message = ""
    if not MOVIEPY_AVAILABLE:
        transcribe_deps_message = """
        <div class="bg-red-100 border-l-4 border-red-500 text-red-700 p-3 mb-4 rounded-lg shadow-md" role="alert">
            <p class="font-bold">Local Video Transcription Warning</p>
            <p class="text-sm">Local file video transcription requires <code>moviepy</code>, <code>pydub</code>, and audio system libraries. Use `pip install moviepy pydub`.</p>
        </div>
        """
    
    pytube_status_message = ""
    if not PYTUBE_AVAILABLE:
        pytube_status_message = """
        <div class="bg-red-100 border-l-4 border-red-500 text-red-700 p-3 mb-4 rounded-lg shadow-md" role="alert">
            <p class="font-bold">YouTube Transcription Disabled</p>
            <p class="text-sm">YouTube transcription requires the **pytube** library. Please install it with: <code>pip install pytube</code></p>
        </div>
        """
    
    
    html_content = f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>📚 Dynamic Voice Assistant</title>
    <script src="https://cdn.tailwindcss.com"></script>
    <script>
        tailwind.config = {{
            theme: {{
                extend: {{
                    fontFamily: {{
                        sans: ['Inter', 'sans-serif'],
                    }},
                    colors: {{
                        'primary': '#10b981', /* Emerald Green */
                        'primary-dark': '#059669',
                        'secondary': '#3b82f6', /* Blue */
                        'surface': '#ffffff',
                    }}
                }}
            }}
        }}
    </script>
    <style>
        /* Custom styles for better scroll and text display */
        #output-label, #file-list-container {{
            min-height: 150px;
            max-height: 400px;
            overflow-y: auto;
            white-space: pre-wrap;
            word-wrap: break-word;
            border-radius: 8px; 
        }}
        .container {{
            border-radius: 12px;
        }}
        .file-item:hover {{
            background-color: #f0fdf4; /* Light green hover */
        }}
        .file-group .file-item:last-child {{
            border-bottom: none !important;
        }}
    </style>
</head>
<body class="bg-gray-50 font-sans p-4 sm:p-8">
    <div class="max-w-6xl mx-auto">
        <h1 class="text-3xl font-extrabold text-primary text-center mb-6">
            <svg xmlns="http://www.w3.org/2000/svg" class="h-8 w-8 inline-block mr-2 text-secondary" fill="none" viewBox="0 0 24 24" stroke="currentColor">
                <path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M19 11H5m14 0a2 2 0 012 2v6a2 2 0 01-2 2H5a2 2 0 01-2-2v-6a2 2 0 012-2m14 0V9a2 2 0 00-2-2M5 11V9a2 2 0 012-2m0 0V5a2 2 0 012-2h6a2 2 0 012 2v2M7 7h10" />
            </svg>
            Dynamic Media Assistant
        </h1>
            {pytube_status_message}
            {transcribe_deps_message}
            {youtube_status_message}
        <p class="text-center text-sm text-gray-600 mb-8">
            Voice and text command interface for file reading, translation, video transcription, and **YouTube playback/transcription** (e.g., 'transcribe cat video on youtube').
        </p>

        <div class="grid grid-cols-1 lg:grid-cols-3 gap-8">
            
            <div class="lg:col-span-2 bg-surface shadow-xl container p-6 sm:p-8 space-y-6">
                <h2 class="text-xl font-bold text-gray-800 border-b pb-2 mb-4 flex items-center">
                    <svg xmlns="http://www.w3.org/2000/svg" class="h-6 w-6 mr-2 text-primary" fill="none" viewBox="0 0 24 24" stroke="currentColor"><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M11 15v2m0 0v2m0-2h2m-2 0H9m-3 4h12a2 2 0 002-2v-6a2 2 0 00-2-2H6a2 2 0 00-2 2v6a2 2 0 002 2zM3 7v2m0 0V7m0 2h2m-2 0H1m8-5h6a2 2 0 012 2v2a2 2 0 01-2 2H9a2 2 0 01-2-2V4a2 2 0 012-2z"/></svg>
                    Assistant Control
                </h2>

                <div>
                    <label for="input-entry" class="block text-sm font-medium text-gray-700 mb-2">
                        Command Input:
                    </label>
                    <input type="text" id="input-entry" placeholder="e.g. read document.pdf, transcribe cat video from Hindi, or reset"
                        class="w-full p-3 border border-gray-300 rounded-lg focus:ring-primary focus:border-primary transition duration-150 shadow-sm" />
                </div>

                <div class="flex flex-col sm:flex-row justify-between space-y-3 sm:space-y-0 sm:space-x-4">
                    <button id="submit-button" onclick="submitCommand()"
                        class="w-full sm:w-1/3 bg-primary text-white font-semibold py-3 px-6 shadow-md hover:bg-primary-dark transition duration-200 rounded-lg">
                        <svg xmlns="http://www.w3.org/2000/svg" class="h-5 w-5 inline-block mr-2" viewBox="0 0 20 20" fill="currentColor"><path d="M17.414 2.586a2 2 0 00-2.828 0L7 10.172V13h2.828l7.586-7.586a2 2 0 000-2.828z" /><path fill-rule="evenodd" d="M2 6a2 2 0 012-2h4a1 1 0 010 2H4v10h10v-4a1 1 0 112 0v4a2 2 0 01-2 2H4a2 2 0 01-2-2V6z" clip-rule="evenodd" /></svg>
                        Submit Text
                    </button>
                    <button id="voice-button" onclick="startVoiceCommand()"
                        class="w-full sm:w-1/3 bg-secondary text-white font-semibold py-3 px-6 shadow-md hover:bg-blue-600 transition duration-200 rounded-lg">
                        <svg xmlns="http://www.w3.org/2000/svg" class="h-5 w-5 inline-block mr-2" viewBox="0 0 20 20" fill="currentColor"><path fill-rule="evenodd" d="M7 4a3 3 0 016 0v4a3 3 0 11-6 0V4z" clip-rule="evenodd" /><path d="M5.5 8A6.5 6.5 0 0112 14.5V16h2a1 1 0 110 2h-4a1 1 0 01-1-1v-1.5A6.501 6.501 0 015 11c0-1.873.79-3.57 2.083-4.793a.75.75 0 00-.75-.025A8.005 8.005 0 004 11a8 8 0 008 8 8.005 8.005 0 004.75-1.782.75.75 0 00-.75-.025A6.501 6.501 0 0114 14.5V16h2a1 1 0 100-2h-2V4a.5.5 0 00-.5-.5z" clip-rule="evenodd" /></svg>
                        Speak Command
                    </button>
                    <button id="reset-button" onclick="submitCommand('forget all')"
                        class="w-full sm:w-1/3 bg-gray-600 text-white font-semibold py-3 px-6 shadow-md hover:bg-gray-700 transition duration-200 rounded-lg">
                        <svg xmlns="http://www.w3.org/2000/svg" class="h-5 w-5 inline-block mr-2" fill="none" viewBox="0 0 24 24" stroke="currentColor" stroke-width="2"><path stroke-linecap="round" stroke-linejoin="round" d="M10 14l2-2m0 0l2-2m-2 2l-2-2m2 2l2 2m7-2a9 9 0 11-18 0 9 9 0 0118 0z" /></svg>
                        Forget All / Reset
                    </button>
                    </div>

                <div class="pt-4">
                    <div class="flex justify-between items-center mb-2">
                        <label class="block text-sm font-bold text-gray-700">
                            Reading/Transcription Output:
                        </label>
                        <span id="read-progress" class="text-sm text-primary font-semibold"></span>
                    </div>
                    <div id="output-label" class="bg-gray-50 border border-gray-300 p-4 text-gray-800 text-base leading-relaxed shadow-inner">
                        <span class="font-bold">Response:</span> Ready.
                    </div>
                    
                    <div class="mt-4 flex space-x-4">
                        <button id="pause-button" onclick="submitCommand('pause')" class="flex-1 bg-gray-200 text-gray-700 py-2 px-4 text-sm font-medium hover:bg-gray-300 transition duration-200 disabled:opacity-50 rounded-lg" disabled>Pause</button>
                        <button id="resume-button" onclick="submitCommand('resume')" class="flex-1 bg-gray-200 text-gray-700 py-2 px-4 text-sm font-medium hover:bg-gray-300 transition duration-200 disabled:opacity-50 rounded-lg" disabled>Resume</button>
                        <button id="restart-button" onclick="submitCommand('restart')" class="flex-1 bg-gray-200 text-gray-700 py-2 px-4 text-sm font-medium hover:bg-gray-300 transition duration-200 disabled:opacity-50 rounded-lg" disabled>Restart</button>
                    </div>
                </div>
            </div>

            <div class="lg:col-span-1 space-y-8">
                
                <div class="bg-surface shadow-xl container p-6 sm:p-8 space-y-4">
                    <h2 class="text-xl font-bold text-gray-800 border-b pb-2 mb-4 flex items-center">
                        <svg xmlns="http://www.w3.org/2000/svg" class="h-6 w-6 mr-2 text-secondary" fill="none" viewBox="0 0 24 24" stroke="currentColor"><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M3 7v10a2 2 0 002 2h14a2 2 0 002-2V9a2 2 0 00-2-2h-6l-2-2H5a2 2 0 00-2 2z" /></svg>
                        File Explorer (Downloads)
                    </h2>
                    <div id="file-list-container" class="border border-gray-200 bg-gray-50 p-2 text-sm">
                        Loading files...
                    </div>
                    <button onclick="listFiles()" class="w-full bg-gray-100 text-gray-700 py-2 text-sm font-medium border border-gray-300 hover:bg-gray-200 transition duration-150 rounded-lg">
                        Refresh File List
                    </button>
                    <p class="text-xs text-gray-500 mt-2">
                        Files are read from the server's configured storage path: <code>{FILE_STORAGE_PATH}</code>
                    </p>
                </div>

                <div class="bg-surface shadow-xl container p-6 sm:p-8 space-y-4">
                    <h2 class="text-xl font-bold text-gray-800 border-b pb-2 mb-4 flex items-center">
                        <svg xmlns="http://www.w3.org/2000/svg" class="h-6 w-6 mr-2 text-primary" fill="none" viewBox="0 0 24 24" stroke="currentColor"><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M3 5h12M9 3v2m1.066 14.885A18 18 0 0110.5 20h2a18 18 0 01-.566-2.115M21 12a9 9 0 11-18 0 9 9 0 0118 0z" /></svg>
                        Translation Settings
                    </h2>
                    
                    <div class="grid grid-cols-1 md:grid-cols-2 gap-6"> 
                        
                        <div>
                            <label for="lang-select" class="block text-sm font-medium text-gray-700">
                                Target Translation Language:
                            </label>
                            <select id="lang-select" onchange="setTranslationLanguage(this.value)"
                                class="w-full p-3 border border-gray-300 rounded-lg focus:ring-primary focus:border-primary transition duration-150 shadow-sm">
                                {language_options}
                            </select>
                            <p class="text-xs text-gray-500 mt-2">
                                This setting applies to file reading and is the **default target** for transcription.
                            </p>
                        </div>

                        <div>
                            <label for="voice-select" class="block text-sm font-medium text-gray-700">
                                Reading Voice:
                            </label>
                            <select id="voice-select" onchange="setVoice(this.value)"
                                class="w-full p-3 border border-gray-300 rounded-lg focus:ring-primary focus:border-primary transition duration-150 shadow-sm">
                                <option>Loading voices...</option>
                            </select>
                            <p class="text-xs text-gray-500 mt-2">
                                Choose a specific voice for text-to-speech reading.
                            </p>
                        </div>
                    </div>
                    <label for="rate-slider" class="block text-sm font-medium text-gray-700 mt-4">
                        Reading Speed: <span id="current-rate-label" class="font-semibold text-primary">1.0x</span>
                    </label>
                    <input type="range" id="rate-slider" min="0.5" max="2.0" value="1.0" step="0.1" 
                        oninput="setReadingRate(this.value)"
                        class="w-full h-2 bg-gray-200 rounded-lg appearance-none cursor-pointer range-lg">
                    <div class="flex justify-between text-xs text-gray-500">
                        <span>0.5x (Slow)</span>
                        <span>2.0x (Fast)</span>
                    </div>
                    </div>
                
            </div>
            
        </div>
    </div>
    
    <script>
        const inputEntry = document.getElementById('input-entry');
        const outputLabel = document.getElementById('output-label');
        const readProgress = document.getElementById('read-progress');
        const fileListContainer = document.getElementById('file-list-container');
        const langSelect = document.getElementById('lang-select');
        const voiceSelect = document.getElementById('voice-select'); 
        const pauseButton = document.getElementById('pause-button');
        const resumeButton = document.getElementById('resume-button');
        const restartButton = document.getElementById('restart-button');
        const rateSlider = document.getElementById('rate-slider'); 
        const currentRateLabel = document.getElementById('current-rate-label'); 
        
        let isReading = false;
        let hasContent = false;
        let selectedVoice = null; 
        let currentRate = 1.0; 
        
        // --- CLIENT-SIDE TTS/UI ---
        
        function setReadingRate(rate) {{
            const newRate = Math.min(2.0, Math.max(0.5, parseFloat(rate)));
            currentRate = newRate;
            
            rateSlider.value = newRate;
            currentRateLabel.textContent = `${{newRate.toFixed(1)}}x`;
            
            if (newRate !== parseFloat(rate)) {{
                safeUiUpdate(`Reading speed adjusted to limit: ${{newRate.toFixed(1)}}x.`, false);
            }}
        }}

        function safeUiUpdate(text, isFileContent = false) {{
            if (isFileContent) {{
                outputLabel.innerHTML = `<span class="font-bold text-primary-dark">Reading:</span> ${{text}}`;
                outputLabel.scrollTop = outputLabel.scrollHeight; 
            }} else {{
                outputLabel.innerHTML = `<span class="font-bold">Response:</span> ${{text}}`;
                readProgress.textContent = ''; 
                updateControlButtons(false);
            }}
        }}
        
        function updateControlButtons(readingActive) {{
            pauseButton.disabled = !readingActive;
            restartButton.disabled = !hasContent;
            
            if (readingActive) {{
                pauseButton.classList.add('bg-red-500', 'text-white');
                pauseButton.classList.remove('bg-gray-200', 'text-gray-700');
                resumeButton.disabled = true;
            }} else {{
                pauseButton.classList.add('bg-gray-200', 'text-gray-700');
                pauseButton.classList.remove('bg-red-500', 'text-white');
                // If paused and has content, enable resume
                resumeButton.disabled = !(hasContent && !isReading);
            }}
            
            if (hasContent && !readingActive) {{
                resumeButton.classList.remove('bg-gray-200', 'text-gray-700');
                resumeButton.classList.add('bg-secondary', 'text-white');
            }} else {{
                resumeButton.classList.add('bg-gray-200', 'text-gray-700');
                resumeButton.classList.remove('bg-secondary', 'text-white');
            }}
        }}
        
        function speak(text) {{
            window.speechSynthesis.cancel(); 
            const utterance = new SpeechSynthesisUtterance(text);
            
            if (selectedVoice) {{
                utterance.voice = selectedVoice;
            }}
            
            utterance.rate = currentRate;

            return new Promise(resolve => {{
                utterance.onend = () => {{
                    resolve();
                }};
                utterance.onerror = (e) => {{
                    console.error("TTS Error:", e);
                    resolve(); 
                }};
                window.speechSynthesis.speak(utterance);
            }});
        }}
        
        function stopReading() {{
            isReading = false;
            window.speechSynthesis.cancel();
            
            if (hasContent) {{
                safeUiUpdate("Reading paused. Say 'resume' or click 'Resume' to continue.", false);
            }} else {{
                safeUiUpdate("Ready.", false);
            }}
            updateControlButtons(false);
        }}
        
        function playYouTubeVideo(searchUrl, query) {{
            window.open(searchUrl, '_blank');
            safeUiUpdate(`Opening YouTube search for: ${{query}}`, false);
        }}


        // --- VOICE SELECTION LOGIC ---

        function loadVoices() {{
            const populateVoiceList = () => {{
                const voices = window.speechSynthesis.getVoices();
                voiceSelect.innerHTML = ''; 

                if (voices.length === 0) {{
                    voiceSelect.innerHTML = '<option value="">No voices available.</option>';
                    selectedVoice = null;
                    return;
                }}

                const targetLangCode = langSelect.value;
                let defaultVoiceName = '';

                const finalVoices = voices.filter(voice => voice.lang.includes(targetLangCode));
                const voicesToDisplay = finalVoices.length > 0 ? finalVoices : voices;

                voicesToDisplay.forEach((voice, index) => {{
                    const option = document.createElement('option');
                    option.textContent = `${{voice.name}} (${{voice.lang}})`;
                    option.value = voice.name;

                    if (voice.default || index === 0) {{
                        option.selected = true;
                        defaultVoiceName = voice.name;
                    }}

                    voiceSelect.appendChild(option);
                }});
                
                const initialVoice = voices.find(v => v.name === defaultVoiceName);
                if (initialVoice) {{
                    selectedVoice = initialVoice;
                }}
            }};

            if ('onvoiceschanged' in window.speechSynthesis) {{
                window.speechSynthesis.onvoiceschanged = populateVoiceList;
            }} 
            populateVoiceList();
        }}

        function setVoice(voiceName) {{
            const voices = window.speechSynthesis.getVoices();
            const newVoice = voices.find(v => v.name === voiceName);
            
            if (newVoice) {{
                selectedVoice = newVoice;
                safeUiUpdate(`Reading voice set to ${{newVoice.name}}.`);
            }} else {{
                safeUiUpdate("Error: Could not find the selected voice.", false);
            }}
        }}

        // --- FLASK API INTERACTION (Updated setTranslationLanguage) ---

        function setTranslationLanguage(langCode) {{
            const langName = langSelect.options[langSelect.selectedIndex].text;
            
            submitCommand(`translate to ${{langName}}`);
            
            loadVoices(); 
        }}
        
        async function submitCommand(overrideCommand = null) {{
            const command = (overrideCommand || inputEntry.value).trim().toLowerCase();
            
            if (command === 'stop' || command === 'pause') {{
                stopReading();
                return; 
            }}
            
            safeUiUpdate(`Processing command: ${{command}}`);
            stopReading(); 
            
            try {{
                const response = await fetch('/command', {{
                    method: 'POST',
                    headers: {{ 'Content-Type': 'application/json' }},
                    body: JSON.stringify({{ question: command }})
                }});
                
                const data = await response.json();
                const responseText = data.message;
                
                // 1. Handle server action
                if (data.action === 'start_read') {{
                    hasContent = true;
                    startReading();
                }} else if (data.action === 'open_url') {{
                    window.open(data.url, '_blank');
                }} else if (data.action === 'set_rate') {{
                    setReadingRate(data.rate);
                }} else if (data.action === 'stop_read') {{
                    hasContent = false;
                    updateControlButtons(false);
                }} 
                // --- YOUTUBE CLIENT-SIDE HANDLERS ---
                else if (data.action === 'play_youtube_embedded') {{
                    // FIX: Open the direct YouTube watch link
                    window.open(`https://www.youtube.com/watch?v=${{data.video_id}}`, '_blank');
                    safeUiUpdate(`Opening YouTube video: ${{data.video_title}}`, false);
                    await speak(data.message);
                    return; 
                }} else if (data.action === 'play_youtube') {{
                    playYouTubeVideo(data.url, data.query);
                    await speak(data.message);
                    return; 
                }}
                // --- END YOUTUBE HANDLERS ---

                // 2. Update UI and speak the conversational response
                safeUiUpdate(responseText, false);
                await speak(responseText); 
                
                // 3. Update hasContent flag and control buttons
                if (data.hasOwnProperty('has_content')) {{
                    hasContent = data.has_content;
                }}
                updateControlButtons(data.action === 'start_read');
                
                if (!overrideCommand) {{
                    inputEntry.value = '';
                }}

            }} catch (error) {{
                const errorMessage = `Error communicating with server: ${{error.message}}`;
                safeUiUpdate(errorMessage, false);
                speak("I encountered an error communicating with the server.");
                console.error("Fetch Error:", error);
                updateControlButtons(false);
            }}
        }}
        
        // Recursive function to fetch, speak, and request the next chunk
        async function fetchAndSpeakChunk() {{
            if (!isReading) return; 
            
            try {{
                const response = await fetch('/read_chunk');
                const data = await response.json();

                if (data.status === 'reading') {{
                    safeUiUpdate(data.chunk, true);
                    readProgress.textContent = `Progress: ${{data.progress}}%`;
                    
                    await speak(data.chunk);

                    if (isReading) {{
                        setTimeout(fetchAndSpeakChunk, 50); 
                    }}
                }} else if (data.status === 'done') {{
                    stopReading(); 
                    safeUiUpdate(data.message, false);
                    speak(data.message);
                    readProgress.textContent = '100% Complete.';
                    hasContent = false; 
                    updateControlButtons(false);
                }} else {{
                    stopReading(); 
                    safeUiUpdate(data.message, false);
                    speak("An error occurred during reading.");
                }}
            }} catch (error) {{
                console.error("Chunk Fetch Error:", error);
                stopReading();
                safeUiUpdate("Lost connection to server during read.", false);
            }}
        }}

        async function startReading() {{
            if (isReading) return;

            if (!hasContent) {{
                safeUiUpdate("Please load a file first using 'read [filename.ext]'.", false);
                return;
            }}
            
            isReading = true;
            updateControlButtons(true);
            fetchAndSpeakChunk();
        }}

        function toggleFolder(folderId, button) {{
            const folder = document.getElementById(folderId);
            const icon = button.querySelector('.toggle-icon');
            
            if (folder.classList.contains('hidden')) {{
                folder.classList.remove('hidden');
                icon.textContent = '▲';
            }} else {{
                folder.classList.add('hidden');
                icon.textContent = '▼';
            }}
        }}
        
        async function listFiles() {{
            fileListContainer.innerHTML = '<p class="text-center py-4 text-gray-500">Fetching file list...</p>';
            try {{
                const response = await fetch('/list_files');
                const data = await response.json();

                if (data.status === 'success' && data.files.length > 0) {{
                    const selectedLangName = langSelect.options[langSelect.selectedIndex].text;
                    
                    const groupedFiles = data.files.reduce((acc, file) => {{
                        const fileExtension = file.split('.').pop().toLowerCase();
                        const groupName = fileExtension.toUpperCase() + ' Files';
                        if (!acc[groupName]) {{
                            acc[groupName] = [];
                        }}
                        acc[groupName].push(file);
                        return acc;
                    }}, {{}});
                    
                    let fileHtml = '';

                    for (const groupName in groupedFiles) {{
                        const fileList = groupedFiles[groupName];
                        
                        const fileExtension = groupName.split(' ')[0].toLowerCase();
                        
                        const isVideoGroup = fileExtension.match(/mp4|mov|avi|mkv|wmv/i); 
                        let groupIcon = isVideoGroup ? '▶️' : '📄';
                        let defaultCommand = isVideoGroup ? `transcribe  ` : 'read ';

                        const filesInGroupHtml = fileList.map(file => {{
                            return `
                                <div class="file-item flex justify-between items-center px-4 py-3 last:border-b-0 cursor-pointer" 
                                    onclick="inputEntry.value='${{defaultCommand}} ${{file}}'; submitCommand();">
                                    <span class="text-gray-700 w-3/4 truncate">${{groupIcon}} ${{file}}</span>
                                    <button class="text-primary hover:text-primary-dark font-medium text-xs focus:outline-none">
                                        ${{isVideoGroup ? 'Transcribe' : 'Read'}}
                                    </button>
                                </div>
                            `;
                        }}).join('');
                        
                        fileHtml += `
                            <div class="mb-2 border border-gray-200 rounded-lg shadow-sm">
                                <button class="w-full text-left p-3 font-semibold text-gray-800 bg-white hover:bg-gray-100 flex justify-between items-center rounded-t-lg" 
                                        onClick="toggleFolder('folder-${{fileExtension}}', this)">
                                    <span class="text-secondary">📁 ${{groupName}} (${{fileList.length}})</span>
                                    <span class="toggle-icon text-sm">▼</span>
                                </button>
                                <div id="folder-${{fileExtension}}" class="file-group hidden bg-white">
                                    ${{filesInGroupHtml}}
                                </div>
                            </div>
                        `;
                    }}

                    fileListContainer.innerHTML = fileHtml;

                }} else if (data.files.length === 0) {{
                    fileListContainer.innerHTML = '<p class="text-center text-gray-500 py-4">No supported files found in the Downloads folder.</p>';
                }} else {{
                    fileListContainer.innerHTML = `<p class="text-red-500">Error: ${{data.message}}</p>`;
                }}
            }} catch (error) {{
                fileListContainer.innerHTML = `<p class="text-red-500">Failed to connect to server: ${{error.message}}</p>`;
            }}
        }}

        // --- CLIENT-SIDE STT (Voice Command) ---
        function startVoiceCommand() {{
            if (!('webkitSpeechRecognition' in window) && !('SpeechRecognition' in window)) {{
                safeUiUpdate("Web Speech Recognition not supported in this browser. Please use text input.", false);
                speak("Speech recognition is not available.");
                return;
            }}

            safeUiUpdate("Listening...", false);
            const SpeechRecognition = window.SpeechRecognition || window.webkitSpeechRecognition;
            const recognition = new SpeechRecognition();
            
            recognition.interimResults = false;
            recognition.maxAlternatives = 1;
            recognition.lang = 'en-US'; // Default client STT to English

            recognition.onresult = (event) => {{
                const speechResult = event.results[0][0].transcript;
                inputEntry.value = speechResult;
                submitCommand(); 
            }};

            recognition.onerror = (event) => {{
                if (event.error !== 'no-speech') {{
                    safeUiUpdate(`Voice error: ${{event.error}}. Try again.`, false);
                    speak("Sorry, I couldn't process that command.");
                }} else {{
                    safeUiUpdate("Response: Ready.", false); 
                }}
            }};
            
            recognition.onend = () => {{
                if (outputLabel.textContent === "Response: Listening...") {{
                    safeUiUpdate("Response: Ready.", false); 
                }}
            }};

            try {{
                recognition.start();
            }} catch (e) {{
                console.warn("Recognition already started or error in start:", e);
                safeUiUpdate("Response: Recognition already started. Speak now.", false);
            }}
        }}
        
        document.addEventListener('DOMContentLoaded', async () => {{
            safeUiUpdate("Ready. Use the File Explorer or speak a command.", false); 
            listFiles(); 
            loadVoices(); 
            hasContent = false;
            updateControlButtons(false);
            
            const initialLangCode = '{translation_language}';
            langSelect.value = initialLangCode;

            setReadingRate(currentRate);
        }});

    </script>
</body>
</html>
"""
    return render_template_string(html_content)

@app.route('/command', methods=['POST'])
def command_handler():
    """
    Handles all conversational and control commands (read, translate, info, web actions, YOUTUBE).
    """
    global current_file_content

    try:
        data = request.json
        question = data.get('question', '').strip().lower()

        has_content = bool(current_file_content)
        
        # 1. Handle RESET Command
        if question in ("reset", "forget", "clear memory", "forget all"):
            response = handle_reset_command()
            return jsonify(response)
        
        # 2. Handle YouTube Playback Commands
        if ("play" in question or "watch" in question) and "youtube" in question and "transcribe" not in question and "translate" not in question:
            youtube_response = handle_youtube_command(question)
            if youtube_response:
                youtube_response['has_content'] = has_content
                return jsonify(youtube_response)

        # 2.5. Handle YouTube Transcription Commands (NOW REAL)
        if ("transcribe" in question or "translate" in question) and "youtube" in question:
            query = extract_youtube_query(question)
            
            # --- Extract and default both language codes ---
            target_lang, source_lang = extract_language_codes(question)
            target_lang = target_lang if target_lang else translation_language  
            source_lang = source_lang if source_lang else 'en'
            
            if not query:
                response_text = "Please specify the video you want to transcribe (e.g., 'transcribe my song on youtube')."
                return jsonify({"status": "error", "message": response_text, "has_content": has_content})
            
            # --- Pass both language codes to the transcription function ---
            response = transcribe_youtube_audio_real(query, target_lang, source_lang)
            
            # Update has_content based on whether a new script was loaded
            response['has_content'] = bool(response.get('action') == 'start_read')
            return jsonify(response)


        # 3. Handle Web Actions (Order/Search) - These return an 'open_url' action for the client
        web_response = handle_web_action_command(question)
        if web_response:
            web_response['has_content'] = has_content
            return jsonify(web_response)

        # 4. Handle File Reading/Opening/Transcribing
        file_name = extract_filename_from_command(question)
        if file_name and any(cmd in question for cmd in ["open", "read", "transcribe", "play"]):
            # --- Extract source/target for local transcription ---
            target_lang, source_lang = extract_language_codes(question)
            
            # This logic is handled in handle_open_file, which now gets language codes
            response = handle_open_file(file_name, question)
            
            if response['status'] == 'success':
                has_content = True
            
            response['has_content'] = has_content
            return jsonify(response)
        
        # 5. Handle Translation
        if "translate to" in question:
            response = handle_translation_command(question)
            if has_content and response['status'] == 'success':
                response['action'] = 'start_read' # Restart reading to apply new language
            response['has_content'] = has_content
            return jsonify(response)

        # 6. Handle Reading Control (Including new rate control)
        if question in ("restart", "resume", "stop", "pause", "continue") or question.startswith("set reading speed to"):
            if question == "continue": question = "resume" # Map 'continue' to 'resume'
            response = handle_control_command(question)
            response['has_content'] = has_content
            return jsonify(response)

        # 7. Handle General Conversation
        if "your name" in question:
            response_text = "I am Assistant AI, a helpful web assistant."
        elif "my name" in question:
            response_text = "I don't store your personal name, but I'm happy to help you."
        elif "time" in question:
            response_text = f"The current server time is {datetime.now().strftime('%I:%M %p')}."
        elif "date" in question:
            response_text = f"Today's date is {datetime.now().strftime('%Y-%m-%d')}."
        elif "joke" in question:
            response_text = "Why don’t scientists trust atoms? Because they make up everything!"
        elif "weather" in question:
            response_text = "I can't check the weather, but I can search the web for it if you ask me to 'search weather in [city]'."
        else:
            response_text = "I am not sure about that. Please try a file command, a control command, a web search, a YouTube command, or say **'reset'** to clear the current file."
        
        return jsonify({
            "status": "info",
            "message": response_text,
            "has_content": has_content
        })

    except Exception as e:
        print(f"Server Error in /command: {e}")
        return jsonify({"status": "error", "message": f"A server error occurred: {e}", "has_content": False}), 500

@app.route('/read_chunk')
def read_chunk_handler():
    """Endpoint for the client to request the next translated chunk."""
    response = get_next_chunk(chunk_size=500)
    return jsonify(response)

if __name__ == '__main__':
    print(f"--- Voice Assistant Server ---")
    print(f"Place readable files (.txt, .pdf, .docx, .pptx, .mp4, etc.) in the: {FILE_STORAGE_PATH} folder.")
    print(f"YouTube Playback Status: {'ENABLED' if REQUESTS_AVAILABLE else 'DISABLED (install requests)'}")
    print(f"YouTube Transcription Status: {'ENABLED (using pytube)' if PYTUBE_AVAILABLE else 'DISABLED (install pytube)'}")
    print(f"Running on http://127.0.0.1:5000/")
    app.run(debug=True)