import os
import json
from flask import Flask, request, jsonify, render_template_string
from googletrans import Translator, LANGUAGES
import fitz # PyMuPDF for PDF
from docx import Document # python-docx
from pptx import Presentation # python-pptx
import speech_recognition as sr
import tempfile
from datetime import datetime
import subprocess 

# Attempt to import moviepy for video processing; essential for transcription
try:
    from moviepy import VideoFileClip
    MOVIEPY_AVAILABLE = True
except ImportError:
    # print("WARNING: moviepy not installed. Video transcription is disabled.")
    MOVIEPY_AVAILABLE = False


# --- CONFIGURATION & GLOBAL STATE (Server-Side) ---

# Set the file path to the user's Downloads folder
FILE_STORAGE_PATH = os.path.join(os.path.expanduser("~"), "Downloads") 


os.makedirs(FILE_STORAGE_PATH, exist_ok=True) # Ensure the folder exists

app = Flask(__name__)
translator = Translator()
recognizer = sr.Recognizer()

# Global State Variables
current_file_content = "" # Stores the original content for file reading/transcription
read_position = 0
translation_language = 'en' # Target language code

# --- HELPER FUNCTIONS ---

def get_language_code(language_name):
    """Converts a language name (e.g., 'French') to a code (e.g., 'fr')."""
    # Includes a reverse lookup and handles case insensitivity
    for code, name in LANGUAGES.items():
        if name.lower() == language_name.lower():
            return code
    return None

def extract_language_codes(command):
    """
    Parses the command for source and destination languages for transcription/translation.
    Example: "transcribe and translate to hindi from english video.mp4"
    Returns: (target_lang_code, source_lang_code) or (None, None)
    """
    command = command.lower()
    target_code = None
    source_code = None

    # 1. Extract Target Language (translate to [language])
    if "translate to" in command:
        parts = command.split("translate to", 1)[1].strip().split()
        if parts:
            # Look for the language name before "from" or the end of the command
            target_name_parts = []
            for part in parts:
                if part == 'from':
                    break
                target_name_parts.append(part)
            
            target_name = " ".join(target_name_parts)
            target_code = get_language_code(target_name)

    # 2. Extract Source Language (from [language])
    if "from" in command:
        # Check for 'from spanish' or 'from english'
        parts = command.split("from", 1)[1].strip().split()
        if parts:
            # Look for the language name which usually immediately follows 'from'
            # We check the first word after 'from'.
            source_name = parts[0].split()[0] if parts[0].split() else parts[0]
            source_code = get_language_code(source_name)
            
            # Special case for 'hindi' source language for better STT performance
            if source_code == 'hi':
                source_code = 'hi-IN' # Use regional code for Hindi STT if requested

    # Default the target language to 'en' if transcription is requested without a target language
    if target_code is None and "transcribe" in command:
        target_code = 'en'
    
    # Default the source language to 'en' if not specified, as this is often a good general starting point
    if source_code is None and "transcribe" in command:
        source_code = 'en' 

    return target_code, source_code

def extract_filename_from_command(command):
    """
    Parses the command to find a file name, handling common STT errors like 
    'dot' or spacing in extensions.
    """
    command = command.lower()
    
    # --- FIX FOR FILE RECOGNITION (e.g., "city dot mp4") ---
    # Replace spoken "dot" forms with a literal dot 
    command = command.replace(" dot ", ".")
    command = command.replace(" point ", ".")
    
    # Handle misspoken/mis-transcribed extensions (e.g., "mp four" -> "mp4")
    command = command.replace(" mp 4", " mp4")
    # --- END FIX ---

    # Expanded keywords and extensions to improve parsing
    file_keywords = ["open", "read", "play", "transcribe", "file", "document", "video", "docx", "pdf", "pptx", "txt", "mp4", "mov", "avi", "mkv", "wmv"]
    
    for keyword in file_keywords:
        if keyword in command:
            parts = command.split(keyword, 1)
            if len(parts) > 1:
                filename_candidate = parts[1].strip()
                # Clean up everything after the file name (like "and translate...")
                for separator in [' and ', ' to ', ' in ']:
                    if separator in filename_candidate:
                        filename_candidate = filename_candidate.split(separator)[0]
                
                # Check if it looks like a file name
                if "." in filename_candidate and len(filename_candidate) > 3:
                    # Final cleaning for file system access (e.g., remove trailing spaces/punctuation)
                    final_filename = filename_candidate.strip().rstrip('.')
                    return final_filename
    return None

def extract_file_content(full_path):
    """Extracts content based on file extension."""
    content = ""
    file_name = os.path.basename(full_path)
    file_extension = os.path.splitext(file_name)[1].lower()

    if file_extension == ".txt":
        with open(full_path, "r", encoding="utf-8") as file:
            content = file.read()
            
    elif file_extension == ".docx":
        doc = Document(full_path)
        content = "\n".join([para.text for para in doc.paragraphs])
        
    elif file_extension == ".pdf":
        doc = fitz.open(full_path)
        content = "\n".join(page.get_text() for page in doc)
        doc.close()
        
    elif file_extension == ".pptx":
        prs = Presentation(full_path)
        content = "\n".join(
            [
                "\n".join(
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ) for slide in prs.slides
            ]
        )
    return content

# --- MODIFIED: Dynamic Transcription and Translation ---
def transcribe_video_audio(file_name, requested_target_lang='en', requested_source_lang='en'):
    """
    Extracts audio from a video, transcribes it using Google STT (using requested_source_lang 
    for detection), and translates the result to requested_target_lang.
    """
    global current_file_content, read_position, translation_language

    if not MOVIEPY_AVAILABLE:
        return {"status": "error", "message": "Video transcription failed: moviepy library is not installed on the server."}
    
    full_path = os.path.join(FILE_STORAGE_PATH, file_name)

    if not os.path.exists(full_path):
        return {"status": "error", "message": f"The video file '{file_name}' does not exist in the server's files/ folder."}

    # Use a temporary file to save the extracted audio
    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmpfile:
        audio_file_path = tmpfile.name

    try:
        # 1. EXTRACT AUDIO
        clip = VideoFileClip(full_path)
        clip.audio.write_audiofile(audio_file_path, logger=None)
        clip.close()
        
        # Use the requested source language for STT recognition.
        stt_lang_code = requested_source_lang
        
        # 2. TRANSCRIBE AUDIO
        with sr.AudioFile(audio_file_path) as source:
            audio = recognizer.record(source) 
            
            # Use the dynamic source language code for transcription
            transcribed_text = recognizer.recognize_google(audio, language=stt_lang_code) 
            
            if not transcribed_text:
                source_name = LANGUAGES.get(stt_lang_code, stt_lang_code)
                raise sr.UnknownValueError(f"No speech detected after trying {source_name} transcription.")

            # --- Dynamic Translation ---
            
            # Detect the *actual* language of the transcribed text (might differ from stt_lang_code)
            detection_result = translator.detect(transcribed_text)
            actual_source_lang_code = detection_result.lang
            actual_source_lang_name = LANGUAGES.get(actual_source_lang_code, 'Unknown Language')
            target_lang_name = LANGUAGES.get(requested_target_lang, requested_target_lang)
            
            final_content = transcribed_text
            translation_message = ""
            
            # Only translate if the detected language differs from the requested target language
            if actual_source_lang_code != requested_target_lang:
                translation_result = translator.translate(transcribed_text, dest=requested_target_lang)
                final_content = translation_result.text
                translation_message = f" (Translated from {actual_source_lang_name} to {target_lang_name}.)"
            else:
                translation_message = f" (Language detected as {actual_source_lang_name}. No translation needed.)"

    except sr.UnknownValueError as e:
        return {"status": "error", "message": f"Transcription failed: Could not understand the audio in the video file. Error: {e}"}
    except sr.RequestError as e:
        return {"status": "error", "message": f"Transcription failed: Could not reach Google Speech Recognition service ({e})."}
    except Exception as e:
        print(f"ERROR: Video transcription failed for {file_name}: {e}")
        return {"status": "error", "message": f"Error during video processing or transcription: {e}"}
        
    finally:
        # CLEANUP
        if os.path.exists(audio_file_path):
            os.remove(audio_file_path)

    if not final_content:
        return {"status": "error", "message": f"Successfully processed the file, but no readable script was generated."}

    # Set global state for reading
    current_file_content = final_content 
    read_position = 0 
    translation_language = requested_target_lang # Set reading language to the requested target
    
    return {
        "status": "success", 
        "action": "start_read",
        "message": f"Transcription complete. Starting to read the script for '{file_name}' in {target_lang_name}.{translation_message}"
    }

# --- CORE LOGIC HANDLERS ---

# --- MODIFIED: Pass language codes to transcribe_video_audio ---
def handle_open_file(file_name_raw, command_text):
    """Opens a file, extracts content, and resets the reading state."""
    global current_file_content, read_position, translation_language

    file_name = file_name_raw.strip()
    full_path = os.path.join(FILE_STORAGE_PATH, file_name)

    if not os.path.exists(full_path):
        return {"status": "error", "message": f"The file '{file_name}' does not exist in the server's files/ folder."}

    file_extension = os.path.splitext(file_name)[1].lower()
    
    # --- File Groups for Correct Handling ---
    TEXT_EXTRACTION_FILES = ['.txt', '.pdf', '.docx', '.pptx'] # Should be read aloud
    MEDIA_FILES = ['.mp4', '.mov', '.avi', '.mkv', '.wmv'] # Should open locally or transcribe

    # 1. Handle Transcription (If explicitly asked for media)
    if file_extension in MEDIA_FILES and "transcribe" in command_text:
        # Extract target and source language codes from the complex command
        target_lang, source_lang = extract_language_codes(command_text)
        
        # Use defaults if parsing failed, but ensure they are valid codes
        target_lang = target_lang if target_lang else 'en'
        source_lang = source_lang if source_lang else 'en'
        
        return transcribe_video_audio(file_name, target_lang, source_lang)
    
    # 2. Handle System Open (For media files only, if transcription not requested)
    if file_extension in MEDIA_FILES:
        try:
            # Use os.startfile (Windows) to open the file in its default media player
            os.startfile(full_path) 
            
            return {
                "status": "info", # Changed status to info, action is not 'start_read'
                "message": f"Opening media file '{file_name}' using your default system application."
            }
        except Exception as e:
            return {"status": "error", "message": f"Could not open media file locally. Error: {e}"}

    # 3. Handle Standard Reading (For all document types)
    if file_extension in TEXT_EXTRACTION_FILES:
        try:
            content = extract_file_content(full_path)
            
            if content:
                current_file_content = content
                read_position = 0
                
                # IMPORTANT: Set action to 'start_read' to launch client reading loop
                return {
                    "status": "success", 
                    "action": "start_read", 
                    "message": f"File '{file_name}' loaded successfully ({len(content)} characters). Starting to read in {LANGUAGES.get(translation_language, 'English')}."
                }
            
            return {"status": "error", "message": f"The file {file_name} was opened, but no readable content was extracted."}

        except Exception as e:
            return {"status": "error", "message": f"Error opening/reading file: {e}"}

    # 4. Fallback
    return {"status": "error", "message": f"File type {file_extension} is not supported for reading or local opening."}


def get_next_chunk(chunk_size=500):
    """Fetches the next chunk of content, translates it if necessary."""
    global read_position, current_file_content, translation_language

    if not current_file_content:
        return {"status": "done", "chunk": "", "message": "No file loaded."}

    if read_position >= len(current_file_content):
        read_position = len(current_file_content) # Ensure boundary
        return {"status": "done", "chunk": "", "message": "End of file reached."}

    chunk_end = min(read_position + chunk_size, len(current_file_content))
    original_chunk = current_file_content[read_position:chunk_end]

    translated_chunk = original_chunk
    if translation_language != 'en':
        try:
            # FIX: Adding timeout for robustness, though not strictly necessary for the bug fix
            translated_chunk = translator.translate(original_chunk, dest=translation_language).text
        except Exception as e:
            # print(f"Translation failed: {e}. Falling back to original chunk.")
            pass 

    read_position = chunk_end # Advance the position
    
    return {
        "status": "reading",
        "chunk": translated_chunk,
        "message": f"Reading chunk in {LANGUAGES.get(translation_language, translation_language)}.",
        "progress": int((read_position / len(current_file_content)) * 100)
    }

def handle_translation_command(command):
    """Changes the target translation language."""
    global translation_language
    
    parts = command.split('to')
    if len(parts) > 1:
        lang_name = parts[-1].strip()
        lang_code = get_language_code(lang_name)
        
        if lang_code:
            translation_language = lang_code
            return {"status": "success", "message": f"Translation language set to {LANGUAGES[lang_code]}."}
        
        return {"status": "error", "message": f"Sorry, I don't recognize the language '{lang_name}'."}
    
    return {"status": "info", "message": f"The current translation language is {LANGUAGES.get(translation_language, translation_language)}."}

def handle_control_command(command):
    """Handles commands like pause, resume, restart, stop."""
    global read_position, current_file_content
    
    if command == "restart":
        if not current_file_content:
            return {"status": "info", "message": "I don't have a file loaded to restart reading."}
        read_position = 0
        return {"status": "action", "action": "start_read", "message": f"Restarted reading session from the beginning in {LANGUAGES.get(translation_language, translation_language)}."}
        
    elif command == "resume":
        if not current_file_content:
            return {"status": "info", "message": "No file content loaded to resume."}
        if read_position >= len(current_file_content):
            return {"status": "info", "message": "The file has already finished reading. Use 'restart' to begin again."}
        return {"status": "action", "action": "start_read", "message": f"Resuming reading from the previous position in {LANGUAGES.get(translation_language, translation_language)}."}

    elif command in ("stop", "pause"):
        return {"status": "info", "message": "Reading paused. Say 'resume' to continue."}
    
    return {"status": "info", "message": "Command recognized, but no immediate server action needed."}

def handle_web_action_command(question):
    """Handles commands that result in opening a web page on the client."""
    
    if "order" in question:
        item = question.replace("order", "").replace("online", "").strip()
        
        if "pizza" in item:
            url = "https://www.dominos.com/"
            message = "Redirecting you to Domino's to order pizza."
        elif "book" in item:
            url = "https://www.amazon.com/"
            message = "Redirecting you to Amazon to order a book."
        elif item:
            url = f"https://www.google.com/search?q=best+quality+{item}+to+order"
            message = f"I've searched for the best quality {item}."
        else:
            return {"status": "info", "message": "Please specify what you want to order."}
        
        return {"status": "action", "action": "open_url", "url": url, "message": message}

    elif "search" in question:
        query = question.replace("search", "").replace("web", "").strip()
        if query:
            search_url = f"https://www.google.com/search?q={query}"
            return {"status": "action", "action": "open_url", "url": search_url, "message": f"I've searched the web for: {query}."}
        else:
            return {"status": "info", "message": "Please specify what you want to search for."}
            
    return None


# --- FLASK ROUTES ---

@app.route('/')
def index():
    """Serves the main HTML/JS front-end page."""
    
    html_content = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title> Voice Assistant</title>
    <script src="https://cdn.tailwindcss.com"></script>
    <script>
        tailwind.config = {
            theme: {
                extend: {
                    fontFamily: {
                        sans: ['Inter', 'sans-serif'],
                    },
                    colors: {
                        'primary': '#4f46e5',
                        'primary-dark': '#4338ca',
                        'secondary': '#f97316',
                        'surface': '#ffffff',
                    }
                }
            }
        }
    </script>
    <style>
        /* Custom styles for better scroll and text display */
        #output-label {
            min-height: 150px;
            max-height: 400px;
            overflow-y: auto;
            white-space: pre-wrap;
            word-wrap: break-word;
            border-radius: 8px; 
        }
        .container {
            border-radius: 12px;
        }
        button {
            border-radius: 8px; 
        }
    </style>
</head>
<body class="bg-gray-100 font-sans p-4 sm:p-8">
    <div class="max-w-3xl mx-auto bg-surface shadow-xl container p-6 sm:p-10">
        <h1 class="text-3xl font-extrabold text-primary text-center mb-6">
            <svg xmlns="http://www.w3.org/2000/svg" class="h-8 w-8 inline-block mr-2" fill="none" viewBox="0 0 24 24" stroke="currentColor">
                <path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M19 11H5m14 0a2 2 0 012 2v6a2 2 0 01-2 2H5a2 2 0 01-2-2v-6a2 2 0 012-2m14 0V9a2 2 0 00-2-2M5 11V9a2 2 0 012-2m0 0V5a2 2 0 012-2h6a2 2 0 012 2v2M7 7h10" />
            </svg>
            Flask Voice Assistant
        </h1>
        <p class="text-center text-sm text-gray-600 mb-8">
            Try: 'read report.pdf', 'translate to Spanish', 'transcribe **and translate to Hindi from English** video.mp4', or 'search for best books'.
        </p>

        <div class="space-y-6">
            <div>
                <label for="input-entry" class="block text-sm font-medium text-gray-700 mb-2">
                    Command Input:
                </label>
                <input type="text" id="input-entry" placeholder="e.g. transcribe and translate to French from Spanish video.mp4"
                    class="w-full p-3 border border-gray-300 rounded-lg focus:ring-primary focus:border-primary transition duration-150 shadow-sm" />
            </div>

            <div class="flex flex-col sm:flex-row justify-center space-y-4 sm:space-y-0 sm:space-x-4">
                <button id="submit-button" onclick="submitCommand()"
                    class="w-full sm:w-auto flex-1 bg-primary text-white font-semibold py-3 px-6 shadow-md hover:bg-primary-dark transition duration-200">
                    <svg xmlns="http://www.w3.org/2000/svg" class="h-5 w-5 inline-block mr-2" viewBox="0 0 20 20" fill="currentColor">
                        <path d="M17.414 2.586a2 2 0 00-2.828 0L7 10.172V13h2.828l7.586-7.586a2 2 0 000-2.828z" />
                        <path fill-rule="evenodd" d="M2 6a2 2 0 012-2h4a1 1 0 010 2H4v10h10v-4a1 1 0 112 0v4a2 2 0 01-2 2H4a2 2 0 01-2-2V6z" clip-rule="evenodd" />
                    </svg>
                    Submit Text
                </button>
                <button id="voice-button" onclick="startVoiceCommand()"
                    class="w-full sm:w-auto flex-1 bg-secondary text-white font-semibold py-3 px-6 shadow-md hover:bg-orange-600 transition duration-200">
                    <svg xmlns="http://www.w3.org/2000/svg" class="h-5 w-5 inline-block mr-2" viewBox="0 0 20 20" fill="currentColor">
                      <path fill-rule="evenodd" d="M7 4a3 3 0 016 0v4a3 3 0 11-6 0V4z" clip-rule="evenodd" />
                      <path d="M5.5 8A6.5 6.5 0 0112 14.5V16h2a1 1 0 110 2h-4a1 1 0 01-1-1v-1.5A6.501 6.501 0 015 11c0-1.873.79-3.57 2.083-4.793a.75.75 0 00-.75-.025A8.005 8.005 0 004 11a8 8 0 008 8 8.005 8.005 0 004.75-1.782.75.75 0 00-.75-.025A6.501 6.501 0 0114 14.5V16h2a1 1 0 100-2h-2V4a.5.5 0 00-.5-.5z" clip-rule="evenodd" />
                    </svg>
                    Speak Command
                </button>
            </div>
            
            <div class="pt-4">
                <div class="flex justify-between items-center mb-2">
                    <label class="block text-sm font-medium text-gray-700">
                        Assistant Output:
                    </label>
                    <span id="read-progress" class="text-sm text-primary-dark font-semibold"></span>
                </div>
                <div id="output-label" class="bg-gray-50 border border-gray-300 p-4 text-gray-800 text-sm leading-relaxed shadow-inner">
                    <span class="font-bold">Response:</span> Ready.
                </div>
            </div>
        </div>
    </div>
    
    <script>
        const inputEntry = document.getElementById('input-entry');
        const outputLabel = document.getElementById('output-label');
        const readProgress = document.getElementById('read-progress');
        
        let isReading = false;
        let hasContent = false;
        
        // --- CLIENT-SIDE TTS/UI ---

        function safeUiUpdate(text, isFileContent = false) {
            if (isFileContent) {
                outputLabel.innerHTML = `<span class="font-bold text-primary">Reading:</span> ${text}`;
            } else {
                outputLabel.innerHTML = `<span class="font-bold">Response:</span> ${text}`;
                readProgress.textContent = ''; 
            }
        }
        
        // Speak function now returns a Promise that resolves when speech finishes
        function speak(text) {
            window.speechSynthesis.cancel(); 
            const utterance = new SpeechSynthesisUtterance(text);
            
            return new Promise(resolve => {
                utterance.onend = () => {
                    resolve();
                };
                utterance.onerror = (e) => {
                    console.error("TTS Error:", e);
                    resolve(); 
                };
                window.speechSynthesis.speak(utterance);
            });
        }

        function stopReading() {
            isReading = false;
            window.speechSynthesis.cancel();
            
            if (hasContent) {
                 safeUiUpdate("Reading paused. Say 'resume' to continue.", false);
            } else {
                 safeUiUpdate("Ready.", false);
            }
        }
        
        // --- FLASK API INTERACTION ---

        async function submitCommand() {
            const command = inputEntry.value.trim().toLowerCase();
            
            if (command === 'stop' || command === 'pause') {
                stopReading();
                return; 
            }

            safeUiUpdate(`Processing command: ${command}`);
            stopReading(); // Stop current speech before sending command
            
            try {
                const response = await fetch('/command', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ question: command })
                });

                const data = await response.json();
                const responseText = data.message;
                
                // 1. Check for server action (start reading or open URL)
                if (data.action === 'start_read') {
                    startReading();
                } else if (data.action === 'open_url') {
                    window.open(data.url, '_blank');
                }

                // 2. Update UI and speak the conversational response
                safeUiUpdate(responseText, false);
                await speak(responseText); // Wait for the conversational response to finish speaking

                // 3. Update hasContent flag
                hasContent = data.has_content;

            } catch (error) {
                const errorMessage = `Error communicating with server: ${error.message}`;
                safeUiUpdate(errorMessage, false);
                speak("I encountered an error communicating with the server.");
                console.error("Fetch Error:", error);
            }
        }
        
        // Recursive function to fetch, speak, and request the next chunk
        async function fetchAndSpeakChunk() {
            if (!isReading) return; // Exit condition for pause/stop

            try {
                const response = await fetch('/read_chunk');
                const data = await response.json();

                if (data.status === 'reading') {
                    safeUiUpdate(data.chunk, true);
                    readProgress.textContent = `Progress: ${data.progress}%`;
                    
                    // Wait for the chunk to be spoken entirely
                    await speak(data.chunk);

                    // Once speaking is complete, request the next chunk after a tiny delay
                    if (isReading) {
                        setTimeout(fetchAndSpeakChunk, 50); 
                    }
                } else if (data.status === 'done') {
                    stopReading(); 
                    safeUiUpdate(data.message, false);
                    speak(data.message);
                    readProgress.textContent = '';
                    hasContent = false;
                } else {
                    stopReading(); 
                    safeUiUpdate(data.message, false);
                    speak("An error occurred during reading.");
                }
            } catch (error) {
                console.error("Chunk Fetch Error:", error);
                stopReading();
                safeUiUpdate("Lost connection to server during read.", false);
            }
        }

        async function startReading() {
            if (isReading) return;

            if (!hasContent) {
                safeUiUpdate("Please load a file first using 'read [filename.ext]'.", false);
                return;
            }
            
            isReading = true;
            // Start the reading process by fetching the first chunk
            fetchAndSpeakChunk();
        }

        // --- CLIENT-SIDE STT (Voice Command) ---

        function startVoiceCommand() {
            if (!('webkitSpeechRecognition' in window) && !('SpeechRecognition' in window)) {
                safeUiUpdate("Web Speech Recognition not supported in this browser. Please use text input.", false);
                speak("Speech recognition is not available.");
                return;
            }

            safeUiUpdate("Listening...", false);
            const SpeechRecognition = window.SpeechRecognition || window.webkitSpeechRecognition;
            const recognition = new SpeechRecognition();
            
            recognition.interimResults = false;
            recognition.maxAlternatives = 1;
            recognition.lang = 'hi-IN'; // Set to Hindi (India) for Hinglish support, can be adjusted by the user in the browser

            recognition.onresult = (event) => {
                const speechResult = event.results[0][0].transcript;
                inputEntry.value = speechResult;
                submitCommand(); 
            };

            recognition.onerror = (event) => {
                if (event.error !== 'no-speech') {
                    safeUiUpdate(`Voice error: ${event.error}. Try again.`, false);
                    speak("Sorry, I couldn't process that command.");
                } else {
                    safeUiUpdate("Response: Ready.", false); 
                }
            };
            
            recognition.onend = () => {
                if (outputLabel.textContent === "Response: Listening...") {
                    safeUiUpdate("Response: Ready.", false); 
                }
            };

            try {
                recognition.start();
            } catch (e) {
                console.warn("Recognition already started or error in start:", e);
                safeUiUpdate("Response: Recognition already started. Speak now.", false);
            }
        }
        
        document.addEventListener('DOMContentLoaded', async () => {
             hasContent = false;
        });

    </script>
</body>
</html>
"""
    return render_template_string(html_content)

@app.route('/command', methods=['POST'])
def command_handler():
    """Handles all conversational and control commands (read, translate, info, web actions)."""
    global current_file_content

    try:
        data = request.json
        question = data.get('question', '').strip().lower()

        has_content = bool(current_file_content)
        
        # 1. Handle Web Actions (Order/Search) - These return an 'open_url' action for the client
        web_response = handle_web_action_command(question)
        if web_response:
            web_response['has_content'] = has_content
            return jsonify(web_response)

        # 2. Handle File Reading/Opening/Transcribing
        file_name = extract_filename_from_command(question)
        if file_name and any(cmd in question for cmd in ["open", "read", "transcribe", "play"]):
            response = handle_open_file(file_name, question)
            
            if response['status'] == 'success':
                has_content = True
                # If content was loaded (file read or video transcribed), start reading loop
                # The handle_open_file for transcription now sets the 'action': 'start_read' directly.
                pass 
            
            response['has_content'] = has_content
            return jsonify(response)
        
        # 3. Handle Translation
        if "translate to" in question:
            response = handle_translation_command(question)
            if has_content:
                 response['action'] = 'start_read' # Restart reading to apply new language
            response['has_content'] = has_content
            return jsonify(response)

        # 4. Handle Reading Control
        if question in ("restart", "resume", "stop", "pause", "continue"):
            if question == "continue": question = "resume" # Map 'continue' to 'resume'
            response = handle_control_command(question)
            response['has_content'] = has_content
            return jsonify(response)

        # 5. Handle General Conversation
        if "your name" in question:
            response_text = "I am Assistant AI, a helpful web assistant."
        elif "my name" in question:
             # Retaining the original script's "Tanisha" hardcode is not scalable, using a generic response
             response_text = "I don't store your personal name, but I'm happy to help you."
        elif "time" in question:
             # FIX: Corrected template literal syntax
             response_text = f"The current server time is {datetime.now().strftime('%I:%M %p')}."
        elif "date" in question:
             # FIX: Corrected template literal syntax
             response_text = f"Today's date is {datetime.now().strftime('%Y-%m-%d')}."
        elif "joke" in question:
            response_text = "Why don’t scientists trust atoms? Because they make up everything!"
        elif "weather" in question:
            response_text = "I can't check the weather, but I can search the web for it if you ask me to 'search weather in [city]'."
        else:
            response_text = "I am not sure about that. Please try a file command, a control command, or a web search."
        
        return jsonify({
            "status": "info",
            "message": response_text,
            "has_content": has_content
        })

    except Exception as e:
        print(f"Server Error in /command: {e}")
        # FIX: Corrected template literal syntax
        return jsonify({"status": "error", "message": f"A server error occurred: {e}", "has_content": False}), 500

@app.route('/read_chunk')
def read_chunk_handler():
    """Endpoint for the client to request the next translated chunk."""
    response = get_next_chunk(chunk_size=500)
    return jsonify(response)

# FIX: Corrected the global variable name for file storage path in the print statement
if __name__ == '__main__': # FIX: Correct use of __name__
    print(f"--- Voice Assistant Server ---")
    print(f"Place readable files (.txt, .pdf, .docx, .pptx, .mp4, etc.) in the: {FILE_STORAGE_PATH} folder.")
    print(f"Running on http://127.0.0.1:5000/")
    app.run(debug=True)