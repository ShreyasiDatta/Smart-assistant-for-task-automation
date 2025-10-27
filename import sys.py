import speech_recognition as sr
import pyttsx3
import time
import datetime
import os
import subprocess
import webbrowser
import threading
import fitz 
import tkinter as tk
from googletrans import Translator, LANGUAGES

# --- INITIALIZATION ---
tts_engine = pyttsx3.init()
tts_engine_running = False
current_file_content = "" # Stores the original English content for file reading
read_position = 0          
recognizer = sr.Recognizer()
translator = Translator()
downloads_path = os.path.join(os.path.expanduser("~"), "Downloads") 

# New Global State for Translation
translation_language = 'en' # Default language for reading is English

# --- UI HELPER (Runs on main thread) ---
def safe_ui_update(text, is_file_content=False):
    """Safely updates the output label from any thread."""
    
    # We use a separate function call on the main thread (root.after) to update the UI
    def update_label():
        if is_file_content:
            # Display file content being read
            output_label.config(text=f"Reading ({translation_language}): {text}", justify=tk.LEFT, anchor='w')
        else:
            # Display conversational response
            output_label.config(text=f"Response: {text}", justify=tk.LEFT, anchor='w')
            
    root.after(0, update_label)

# --- SPEECH FUNCTIONS ---
def stop_speaking():
    """Stops the TTS engine and reading thread, preserving the read_position."""
    global tts_engine_running
    if tts_engine_running:
        print("DEBUG: Stopping TTS engine for pause/stop.")
        tts_engine_running = False
        tts_engine.stop()
        
def speak(text, block=False, display_text=None):
    """
    Handles Text-to-Speech. 
    block=True for synchronous file reading.
    block=False (default) for asynchronous conversation.
    """
    global tts_engine_running

    if display_text and not block:
        # For conversational responses, update UI immediately
        safe_ui_update(display_text)

    if block:
        # Synchronous execution for reading chunks sequentially
        try:
            tts_engine.say(text)
            tts_engine.runAndWait()
        except Exception as e:
            print(f"ERROR: Synchronous TTS failed: {e}")
    else:
        # Asynchronous execution for quick conversational responses
        stop_speaking()
        
        def run_tts():
            nonlocal text
            tts_engine_running = True
            try:
                tts_engine.say(text)
                tts_engine.runAndWait()
                if tts_engine_running: 
                    tts_engine_running = False
            except Exception as e:
                print(f"ERROR: Asynchronous TTS failed: {e}")
                tts_engine_running = False

        threading.Thread(target=run_tts).start()

def listen():
    with sr.Microphone() as source:
        try:
            recognizer.adjust_for_ambient_noise(source)
            audio = recognizer.listen(source, timeout=5, phrase_time_limit=10) 
            return recognizer.recognize_google(audio)
        except (sr.UnknownValueError, sr.RequestError):
            return None

def translate_to_english(text):
    return translator.translate(text, src='auto', dest='en').text

def get_language_code(language_name):
    """Converts a language name (e.g., 'French') to a code (e.g., 'fr')."""
    for code, name in LANGUAGES.items():
        if name.lower() == language_name.lower():
            return code
    return None

# --- FILE HANDLING ---
def extract_filename_from_command(command):
    words = command.lower().split()
    file_keywords = ["open", "read", "file", "document", "docx", "pdf", "pptx", "txt"]

    for keyword in file_keywords:
        if keyword in words:
            try:
                start_index = words.index(keyword)
                filename_parts = words[start_index + 1:]
                filename = " ".join(filename_parts)
                if "." in filename and len(filename) > 3:
                    return filename
            except ValueError:
                pass
                
    stripped_command = command.lower().replace("open", "").replace("read", "").strip()
    if "." in stripped_command:
        return stripped_command

    return None

def open_file(file_name_raw):
    global current_file_content, read_position, translation_language
    
    file_name = file_name_raw.strip()
    full_path = os.path.join(downloads_path, file_name) 

    if not os.path.exists(full_path):
        return f"The file '{file_name}' does not exist in your Downloads folder."
    
    try:
        subprocess.Popen(["start", full_path], shell=True) 
        content = ""

        # --- Content Extraction ---
        if file_name.endswith(".txt"):
            with open(full_path, "r", encoding="utf-8") as file:
                content = file.read()
        
        elif file_name.endswith(".docx"):
            from docx import Document
            doc = Document(full_path)
            content = "\n".join([para.text for para in doc.paragraphs])
        
        elif file_name.endswith(".pdf"):
            doc = fitz.open(full_path)
            content = "\n".join(page.get_text() for page in doc)
            doc.close()
            
        elif file_name.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(full_path)
            content = "\n".join(
                [
                    "\n".join(
                        shape.text.strip()
                        for shape in slide.shapes
                        if hasattr(shape, "text") and shape.text.strip()
                    ) for slide in prs.slides
                ]
            )
        
        if content:
            current_file_content = content 
            read_position = 0 
            translation_language = 'en' # Reset to English when opening new file
            
            threading.Thread(target=read_file_content).start() 
            return f"I've opened and started reading the file in English. Say 'translate to [language]' to change."

        return f"The file {file_name} has been opened, but no readable content was extracted."

    except Exception as e:
        return f"Error opening file: {e}"
    

def read_file_content():
    """Reads the file content, translates it, and updates the UI."""
    global read_position, tts_engine_running, current_file_content, translation_language
    
    stop_speaking()
    tts_engine_running = True 
    print(f"DEBUG: Starting file read from position {read_position} in language {translation_language}.")

    while read_position < len(current_file_content):
        if not tts_engine_running:
            print("DEBUG: Reading thread gracefully stopped.")
            break 

        chunk_end = min(read_position + 500, len(current_file_content))
        original_chunk = current_file_content[read_position:chunk_end]
        
        # --- TRANSLATION LOGIC ---
        if translation_language != 'en':
            try:
                translated_chunk = translator.translate(original_chunk, dest=translation_language).text
            except Exception as e:
                print(f"Translation failed: {e}. Falling back to English.")
                translated_chunk = original_chunk
        else:
            translated_chunk = original_chunk
        # --- END TRANSLATION LOGIC ---
        
        # Write the chunk to the UI before speaking it
        safe_ui_update(translated_chunk, is_file_content=True)
        
        # Synchronous speak (blocks until done)
        speak(translated_chunk, block=True) 
        
        read_position = chunk_end 

    if tts_engine_running: 
        tts_engine_running = False
        speak("End of file reached.", display_text="End of file reached.")

# Function to START from the beginning
def restart_reading():
    global read_position, current_file_content, tts_engine_running
    
    if not current_file_content:
        return "I don't have a file loaded to restart reading."
        
    stop_speaking()
    read_position = 0 
    tts_engine_running = True
    threading.Thread(target=read_file_content).start() 
    return f"Restarted reading session from the beginning in {LANGUAGES.get(translation_language, translation_language)}."


# Function to RESUME from the pause point
def resume_reading():
    """Resumes reading the current file content from the last stopped position."""
    global read_position, current_file_content, tts_engine_running
    
    if not current_file_content:
        return "No file content loaded to resume."
        
    if read_position >= len(current_file_content):
        return "The file has already finished reading. Use 'restart' to begin again."
    
    if tts_engine_running:
        return "The file is already being read."

    stop_speaking()
    tts_engine_running = True
    threading.Thread(target=read_file_content).start()
    return f"Resuming reading from the previous position in {LANGUAGES.get(translation_language, translation_language)}."

# --- COMMAND HANDLING ---
def order_online(query):
    if "pizza" in query.lower():
        webbrowser.open("https://www.dominos.com/")
        return "Redirecting you to Domino's."
    elif "book" in query.lower():
        webbrowser.open("https://www.amazon.com/")
        return "Redirecting you to Amazon."
    else:
        # WARNING: Retained bug where this relies on input_entry.get()
        item = input_entry.get().strip() 
        if item:
            search_url = f"https://www.google.com/search?q=best+quality+{item}+to+order"
            webbrowser.open(search_url)
            return f"I've searched for the best quality {item}."
        else:
            return "I need more details."

def search_web(query):
    search_url = f"https://www.google.com/search?q={query}"
    webbrowser.open(search_url)
    return f"I've searched the web for: {query}."

def handle_translation_command(command):
    global translation_language
    
    parts = command.split('to')
    if len(parts) > 1:
        lang_name = parts[-1].strip()
        lang_code = get_language_code(lang_name)
        
        if lang_code:
            translation_language = lang_code
            
            # Restart reading if a file is active
            if tts_engine_running or current_file_content:
                stop_speaking()
                # Do not reset read_position, just restart the reading loop
                threading.Thread(target=read_file_content).start()
                return f"Translating now to {LANGUAGES[lang_code]}. Resuming read."
            
            return f"Set translation language to {LANGUAGES[lang_code]}."
        
        return f"Sorry, I don't recognize the language '{lang_name}'."
    
    return f"The current translation language is {LANGUAGES.get(translation_language, translation_language)}."

def respond(question):
    question = question.lower()
    translated_question = translate_to_english(question)

    # --- Translation Command ---
    if "translate to" in translated_question:
        return handle_translation_command(translated_question)
    
    # --- File/Reading Commands ---
    if translated_question in ("stop", "pause"):
        stop_speaking()
        return "Reading and speaking paused."
    
    elif translated_question in ("resume", "continue"): 
        return resume_reading() 
        
    elif translated_question == "restart":
        return restart_reading()
        
    # --- Other Commands ---
    elif "your name" in translated_question:
        return "My name is AI, your assistant."
    elif "my name" in translated_question:
        return "You are Tanisha, my owner."
    elif "weather" in translated_question:
        return "I cannot check the weather right now."
    elif "time" in translated_question:
        return f"The current time is {datetime.datetime.now().strftime('%I:%M %p')}."
    elif "date" in translated_question:
        return f"Today's date is {datetime.datetime.now().strftime('%Y-%m-%d')}."
    elif "joke" in translated_question:
        return "Why don’t scientists trust atoms? Because they make up everything!"
    elif any(word in translated_question for word in ["open", "read", "file"]):
        file_name = extract_filename_from_command(translated_question)
        return open_file(file_name) if file_name else "Please specify a file name."
    elif "order" in translated_question:
        return order_online(translated_question)
    elif "search" in translated_question:
        return search_web(translated_question)
    
    else:
        return "I am not sure about that."

# --- BUTTON ACTIONS ---
def on_button_click():
    question = input_entry.get().strip().lower()
    answer = respond(question)
    
    # Only speak if the response is not a file action
    if not any(keyword in answer for keyword in ["Reading", "Restarted", "Resuming", "paused", "stopped", "file", "Set translation"]):
        speak(answer, display_text=answer)
    else:
        safe_ui_update(answer)

def on_voice_command():
    safe_ui_update("Listening...")

    question = listen()

    if question:
        safe_ui_update(f"You said: {question}", is_file_content=False)
        
        answer = respond(question)
        
        if not any(keyword in answer for keyword in ["Reading", "Restarted", "Resuming", "paused", "stopped", "file", "Set translation"]):
            speak(answer, display_text=answer)
        else:
            safe_ui_update(answer, is_file_content=False)
            
    else:
        speak("Sorry, I couldn't understand. Please try again.", display_text="Sorry, I couldn't understand. Please try again.")


# --- UI SETUP ---
root = tk.Tk()
root.title("Voice and Text Assistant (with Translation)")
root.geometry("550x450")

label = tk.Label(root, text="Enter command (e.g., 'read report.pdf' or 'translate to Spanish'):")
label.pack(pady=10)

input_entry = tk.Entry(root, width=60)
input_entry.pack(pady=10, padx=10)

button_frame = tk.Frame(root)
button_frame.pack(pady=5)

submit_button = tk.Button(button_frame, text="Submit", command=on_button_click)
submit_button.pack(side=tk.LEFT, padx=10)

voice_button = tk.Button(button_frame, text="Speak", command=on_voice_command)
voice_button.pack(side=tk.LEFT, padx=10)

output_label = tk.Label(root, text="Response: Ready.", wraplength=500, justify=tk.LEFT, anchor='w', borderwidth=2, relief="groove")
output_label.pack(pady=20, fill='x', padx=10, ipady=10)

root.mainloop()